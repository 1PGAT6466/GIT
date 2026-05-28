"""
宝利根知识库后端 — FastAPI + FAISS + Sentence-Transformers
生产级代码，专注安全与健壮性。
"""
from __future__ import annotations

import hashlib
import json
import os
import re
import shutil
import time
import uuid
from datetime import datetime, timezone
from pathlib import Path
try:
    import jieba
    jieba.setLogLevel(20)  # 静默
except ImportError:
    jieba = None
from typing import Optional
import asyncio
import threading

from fastapi import FastAPI, File, HTTPException, Query, Request, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.middleware.gzip import GZipMiddleware
from fastapi.responses import JSONResponse, HTMLResponse, FileResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel, Field
from sentence_transformers import SentenceTransformer
import numpy as np
import faiss

# ============ CONFIG ============
HOST = os.getenv("KB_HOST", "0.0.0.0")
PORT = int(os.getenv("KB_PORT", "8080"))
DB_PATH = os.getenv("KB_DB_PATH", os.path.expanduser("~/kb-server/data"))
UPLOAD_DIR = os.getenv("KB_UPLOAD_DIR", os.path.expanduser("~/kb-server/uploads"))
EMBEDDING_MODEL = os.getenv("KB_MODEL", "BAAI/bge-base-zh-v1.5")
EMBEDDING_DIM = int(os.getenv("KB_EMBEDDING_DIM", "1024"))
CHUNK_SIZE = int(os.getenv("KB_CHUNK_SIZE", "800"))
CHUNK_OVERLAP = int(os.getenv("KB_CHUNK_OVERLAP", "150"))
PARENT_CHUNK_SIZE = int(os.getenv("KB_PARENT_CHUNK_SIZE", "1500"))  # P1: Parent 粒度
CHILD_CHUNK_SIZE = int(os.getenv("KB_CHILD_CHUNK_SIZE", "400"))    # P1: Child 粒度
RERANK_MODEL = os.getenv("KB_RERANK_MODEL", "")  # P2: Qwen/Qwen3-Reranker-0.6B (空=禁用)
RERANK_TOP_K = int(os.getenv("KB_RERANK_TOP_K", "20"))  # 初筛取 top_k 送入 Reranker
BM25_ENABLED = os.getenv("KB_BM25_ENABLED", "true").lower() == "true"  # P2: 混合检索开关
MAX_FILE_MB = int(os.getenv("KB_MAX_FILE_MB", "200"))
MAX_PDF_MB = int(os.getenv("KB_MAX_PDF_MB", "50"))  # PDF 单独限制，大 PDF 解析太慢
ALLOWED_EXTENSIONS = {
    ".txt", ".md", ".csv",
    ".docx", ".doc",
    ".xlsx", ".xls",
    ".pdf",
    ".pptx", ".ppt",
    ".cfg", ".log", ".ini", ".conf", ".json", ".xml", ".html", ".htm",
    ".dxf", ".dwg",
    ".stl", ".obj", ".step", ".stp", ".iges", ".igs", ".gltf", ".glb",
}
BACKUP_DIR = os.getenv("KB_BACKUP_DIR", os.path.expanduser("~/kb-server/backups"))
SEARCH_LOG_DIR = os.getenv("KB_SEARCH_LOG", os.path.expanduser("~/kb-server/logs"))
SEARCH_LOG_MAX = 30
os.makedirs(BACKUP_DIR, exist_ok=True)
os.makedirs(SEARCH_LOG_DIR, exist_ok=True)
SENSITIVE_PATTERNS = [
    re.compile(r"(password|passwd|pwd)\s*[:=]\s*\S+", re.I),
    re.compile(r"(secret|token|api[_-]?key)\s*[:=]\s*\S+", re.I),
    re.compile(r"\b\d{6}(19|20)\d{2}(0[1-9]|1[0-2])(0[1-9]|[12]\d|3[01])\d{3}[\dxX]\b"),  # 身份证
    re.compile(r"\b1[3-9]\d{9}\b"),  # 手机号
]
COLLECTION_NAME = "polygon_knowledge"
START_TIME = time.time()

# ============ APP ============
app = FastAPI(
    title="宝利根知识库 API",
    version="2.0.0",
    docs_url=None,
    redoc_url=None,
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["GET", "POST", "DELETE"],
    allow_headers=["*"],
)
# GZip 压缩：减少 JSON 响应体积 60-80%
app.add_middleware(GZipMiddleware, minimum_size=500)

# Static files
STATIC_DIR = os.path.join(os.path.dirname(__file__), "static")
os.makedirs(STATIC_DIR, exist_ok=True)

@app.get("/", response_class=HTMLResponse)
async def root():
    """首页 — 返回前端 HTML"""
    index_path = os.path.join(STATIC_DIR, "index.html")
    if os.path.exists(index_path):
        with open(index_path, "r", encoding="utf-8") as f:
            return f.read()
    return HTMLResponse("<h1>宝利根知识库 API v2.0</h1><p>请将前端文件放到 static/ 目录</p>")

# 管理面板静态文件
ADMIN_DIR = os.path.join(STATIC_DIR, "admin")
os.makedirs(ADMIN_DIR, exist_ok=True)

# 手动 admin 路由（不用 mount，避免干扰根路径）
@app.get("/admin", response_class=HTMLResponse)
async def admin_home():
    admin_index = os.path.join(ADMIN_DIR, "index.html")
    if os.path.exists(admin_index):
        with open(admin_index, "r", encoding="utf-8") as f:
            return f.read()
    return HTMLResponse("<h1>管理面板未部署</h1>")

@app.get("/admin/{path:path}")
async def admin_files(path: str):
    from fastapi.responses import FileResponse
    import mimetypes
    full_path = os.path.join(ADMIN_DIR, path)
    # 安全：防止目录穿越
    if not os.path.realpath(full_path).startswith(os.path.realpath(ADMIN_DIR)):
        raise HTTPException(status_code=403, detail="Forbidden")
    if os.path.isfile(full_path):
        mime, _ = mimetypes.guess_type(full_path)
        return FileResponse(full_path, media_type=mime)
    # SPA 回退
    admin_index = os.path.join(ADMIN_DIR, "index.html")
    if os.path.exists(admin_index):
        with open(admin_index, "r", encoding="utf-8") as f:
            return HTMLResponse(f.read())
    raise HTTPException(status_code=404, detail="Not found")

# ============ GLOBALS (lazy init) ============
_embedder: Optional[SentenceTransformer] = None
_reranker: Optional[SentenceTransformer] = None  # P2: Cross-Encoder Reranker
_index: Optional[faiss.Index] = None  # FAISS 向量索引
_metadata: list[dict] = []  # 元数据列表，与索引行一一对应
_collection: Optional[str] = None
_ready = False
_meta_lock = threading.Lock()  # 元数据写锁
_index_lock = threading.Lock()  # FAISS 索引写入锁（v3.0+）


# ============ MODELS ============
# 并发上传限制（最多 3 个同时处理）
_upload_semaphore = threading.BoundedSemaphore(3)

class StatsResponse(BaseModel):
    total_files: int
    total_chunks: int
    uptime_seconds: float
    version: str = "2.0.0"
    total_tools: int = 0
    total_faq: int = 0
    feedback_count_today: int = 0


class SearchResult(BaseModel):
    file_name: str
    text_preview: str
    category: str = "未分类"
    tags: list[str] = Field(default_factory=list)
    trust: str = "unverified"
    audit_note: str = ""
    score: float = 0.0


class UploadResponse(BaseModel):
    file: str
    chunks: int
    category: str
    status: str = "ok"


class ChunkPreview(BaseModel):
    index: int
    text: str
    length: int
    flagged: bool = False
    flag_reason: str = ""


class PreviewResponse(BaseModel):
    file_name: str
    file_hash: str
    total_chunks: int
    chunks: list[ChunkPreview]


class ConfirmRequest(BaseModel):
    file_hash: str
    file_name: str
    chunks: list[dict]


class FeedbackRequest(BaseModel):
    """AI 回答反馈"""
    query: str = ""
    answer_preview: str = ""
    useful: bool = True
    timestamp: str = ""


# ============ 可配置数据（工具列表 + 常见问题） ============
# 这些数据原先硬编码在前端 index.html 中，现在集中管理

TOOLS_DATA = [
    {"id": "vpn", "name": "VPN 连接", "icon": "🔒", "category": "网络",
     "url": "#vpn", "desc": "请联系 IT 获取 VPN 客户端与配置", "available": False},
    {"id": "email", "name": "企业邮箱", "icon": "📟", "category": "办公",
     "url": "https://mail.polygon.com", "desc": "宝利根企业邮箱入口", "available": True},
    {"id": "oa", "name": "OA 系统", "icon": "🗚", "category": "办公",
     "url": "#oa", "desc": "审批、考勤、公告", "available": True},
    {"id": "nps", "name": "NPS 认证", "icon": "🛡️", "category": "安全",
     "url": "#nps", "desc": "802.1X 网络准入认证", "available": True},
    {"id": "fileshare", "name": "文件共享", "icon": "📧", "category": "办公",
     "url": "#fileshare", "desc": "部门共享文件夹", "available": True},
    {"id": "printer", "name": "网络打印", "icon": "🖨️", "category": "办公",
     "url": "#printer", "desc": "打印机驱动与配置", "available": True},
    {"id": "phone", "name": "通讯录", "icon": "📓", "category": "办公",
     "url": "#phone", "desc": "全员通讯录查询", "available": True},
    {"id": "ithelp", "name": "IT 工单", "icon": "🎿", "category": "IT",
     "url": "#ithelp", "desc": "报修 / 账号申请", "available": True},
    {"id": "kms", "name": "KMS 激活", "icon": "🔑", "category": "IT",
     "url": "#kms", "desc": "Windows/Office 批量激活", "available": True},
]

FAQ_DATA = [
    {"id": 1, "category": "网络建设",
     "question": "公司网络升级改造方案中，核心交换机和接入层的设备型号是什么？",
     "keywords": ["网络", "交换机", "方案"]},
    {"id": 2, "category": "网络建设",
     "question": "公司内部的 VLAN 划分方案是怎样的？各 VLAN 对应的网段是什么？",
     "keywords": ["VLAN", "网段", "划分"]},
    {"id": 3, "category": "网络建设",
     "question": "网络改造方案的总体预算和分项费用是多少？",
     "keywords": ["预算", "费用", "改造"]},
    {"id": 4, "category": "网络建设",
     "question": "NPS 域控认证的配置流程是怎样的？",
     "keywords": ["NPS", "域控", "认证", "802.1X"]},
    {"id": 5, "category": "IT资产",
     "question": "公司目前有哪些软件授权？各软件的版本和使用部门是什么？",
     "keywords": ["软件", "授权", "版本"]},
    {"id": 6, "category": "IT资产",
     "question": "员工设备的账户规范要求是什么？电脑命名规则是怎样的？",
     "keywords": ["账户", "命名", "规范"]},
    {"id": 7, "category": "IT资产",
     "question": "公司各员工的 IP 地址和对应的电脑信息是什么？",
     "keywords": ["IP", "电脑", "员工"]},
    {"id": 8, "category": "采购合同",
     "question": "网络设备的采购合同中，供货清单包含哪些设备？",
     "keywords": ["采购", "设备", "供货"]},
    {"id": 9, "category": "采购合同",
     "question": "康成报价中综合布线的总金额和明细有哪些？",
     "keywords": ["康成", "布线", "报价"]},
    {"id": 10, "category": "采购合同",
     "question": "办公用品及耗材采购合同的付款方式和交货条款是什么？",
     "keywords": ["付款", "交货", "耗材"]},
    {"id": 11, "category": "公司架构",
     "question": "公司的组织架构是怎样的？各部门的主要职责是什么？",
     "keywords": ["架构", "部门", "职责"]},
    {"id": 12, "category": "自动化",
     "question": "Mini-FakraTE 自动化产线的 CT 要求和 OEE 目标是多少？",
     "keywords": ["CT", "OEE", "产线", "自动化"]},
    {"id": 13, "category": "自动化",
     "question": "Mini-FakraTE 产线适用哪些产品型号？",
     "keywords": ["Mini-Fakra", "产品", "型号"]},
    {"id": 14, "category": "标准件",
     "question": "标准件表中线轨的型号、品牌和规格是什么？",
     "keywords": ["线轨", "标准件", "品牌"]},
]

# ============ HELPERS ============
def _get_embedder() -> SentenceTransformer:
    global _embedder
    if _embedder is None:
        _embedder = SentenceTransformer(EMBEDDING_MODEL)
    return _embedder


def _get_reranker() -> Optional[SentenceTransformer]:
    """P2: Lazy-load Cross-Encoder Reranker"""
    global _reranker
    if not RERANK_MODEL:
        return None
    if _reranker is None:
        try:
            from sentence_transformers import CrossEncoder
            _reranker = CrossEncoder(RERANK_MODEL)
            print(f"[OK] Reranker {RERANK_MODEL} 已加载")
        except Exception as e:
            print(f"[WARN] Reranker 加载失败: {e}")
            return None
    return _reranker


# ============ FAISS 存储层 ============

def _faiss_index_path() -> str:
    return os.path.join(DB_PATH, "index.faiss")

def _faiss_meta_path() -> str:
    return os.path.join(DB_PATH, "metadata.json")


def _load_index():
    """加载 FAISS 索引和元数据"""
    global _index, _metadata
    os.makedirs(DB_PATH, exist_ok=True)
    ipath = _faiss_index_path()
    mpath = _faiss_meta_path()
    if os.path.exists(ipath) and os.path.exists(mpath):
        _index = faiss.read_index(ipath)
        with open(mpath, "r", encoding="utf-8") as f:
            _metadata = json.load(f)
        print(f"[OK] FAISS 索引加载: {_index.ntotal} 条向量")
    else:
        _index = faiss.IndexFlatIP(EMBEDDING_DIM)  # Inner Product = COSINE on normalized vectors
        _metadata = []
        print(f"[OK] FAISS 索引创建 (dim={EMBEDDING_DIM})")


def _save_index():
    """持久化 FAISS 索引和元数据到磁盘（v3.0: 写入加锁防竞态）"""
    with _index_lock:
        if _index is not None:
            faiss.write_index(_index, _faiss_index_path())
    with _meta_lock:
        with open(_faiss_meta_path(), "w", encoding="utf-8") as f:
            json.dump(_metadata, f, ensure_ascii=False)


def _check_dim_compat() -> bool:
    """检查索引维度是否匹配当前模型"""
    if _index is not None and _index.ntotal > 0:
        return _index.d == EMBEDDING_DIM
    return True


def _reset_index():
    """重建空索引"""
    global _index, _metadata
    _index = faiss.IndexFlatIP(EMBEDDING_DIM)
    _metadata = []
    _save_index()


def _faiss_insert(data: list[dict]):
    """批量插入向量到 FAISS 索引"""
    global _index, _metadata
    if not data:
        return
    vectors = np.array([d["vector"] for d in data], dtype=np.float32)
    faiss.normalize_L2(vectors)  # IP → COSINE
    start_idx = _index.ntotal
    _index.add(vectors)
    # 元数据（不存向量，metadata.json 保持轻量）
    with _meta_lock:
        for i, d in enumerate(data):
            meta = {k: v for k, v in d.items() if k != "vector"}
            meta["id"] = start_idx + i
            _metadata.append(meta)
    _save_index()


def _rebuild_faiss_from_meta():
    """从元数据重建 FAISS 索引（删除后）。注意：删除后需业务层重新插入向量。"""
    global _index
    _index = faiss.IndexFlatIP(EMBEDDING_DIM)
    _save_index()


def _sanitize_filename(filename: str) -> str:
    """防路径穿越攻击"""
    name = os.path.basename(filename)
    name = re.sub(r'[<>:"/\\|?*\x00-\x1f]', "_", name)
    if not name or name.startswith("."):
        name = "untitled_" + uuid.uuid4().hex[:8]
    return name


def _extract_text(file_path: str, ext: str) -> str:
    """提取文件纯文本"""
    text = ""
    ext = ext.lower()

    if ext in (".txt", ".md", ".csv", ".cfg", ".log", ".ini", ".conf", ".json", ".xml", ".html", ".htm"):
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()

    elif ext in (".docx", ".doc"):
        try:
            from docx import Document
            doc = Document(file_path)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception:
            text = "[无法解析的 Word 文档]"

    elif ext in (".xlsx", ".xls"):
        try:
            from openpyxl import load_workbook
            wb = load_workbook(file_path, read_only=True, data_only=True)
            parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                parts.append(f"=== {sheet_name} ===")
                for row in ws.iter_rows(values_only=True):
                    parts.append("\t".join(str(c) if c is not None else "" for c in row))
            text = "\n".join(parts)
            wb.close()
        except Exception:
            text = "[无法解析的 Excel 文件]"

    elif ext == ".pdf":
        try:
            import subprocess
            check = subprocess.run(["which", "pdftotext"], capture_output=True, text=True, timeout=5)
            if check.returncode != 0:
                text = "[PDF 解析需要安装 poppler-utils: sudo apt install poppler-utils]"
            else:
                result = subprocess.run(
                    ["pdftotext", "-layout", "-q", file_path, "-"],
                    capture_output=True, text=True, timeout=30
                )
                text = result.stdout.strip()
                if not text or len(text) < 20:
                    text = f"[PDF: {os.path.basename(file_path)} — 无文字内容或为图片型文档]"
        except subprocess.TimeoutExpired:
            text = f"[PDF: {os.path.basename(file_path)} — 解析超时，文件过大]"
        except Exception:
            text = "[无法解析的 PDF 文件]"

    elif ext in (".pptx", ".ppt"):
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
            text = "\n".join(parts)
        except Exception:
            text = "[无法解析的 PPT 文件]"

    elif ext == ".dxf":
        try:
            import ezdxf
            doc = ezdxf.readfile(file_path)
            parts = []
            for entity in doc.modelspace():
                if entity.dxftype() == "TEXT":
                    parts.append(entity.dxf.text)
                elif entity.dxftype() == "MTEXT":
                    parts.append(entity.text)
            text = "\n".join(parts)
            if not text.strip():
                text = f"[DXF: {os.path.basename(file_path)} — 无文字实体]"
        except Exception as e:
            text = f"[DXF: {os.path.basename(file_path)} — 解析失败: {str(e)[:80]}]"

    elif ext == ".dwg":
        try:
            # 方案1: dwgread (LibreDWG, 可能未安装)
            import subprocess
            text = ""
            try:
                result = subprocess.run(
                    ["dwgread", file_path],
                    capture_output=True, text=True, timeout=30
                )
                if result.stdout.strip():
                    text = result.stdout.strip()
            except Exception:
                pass
            # 方案2: strings 命令兜底
            if not text:
                try:
                    result = subprocess.run(
                        ["strings", file_path],
                        capture_output=True, text=True, timeout=30
                    )
                    lines = []
                    for l in result.stdout.split("\n"):
                        s = l.strip()
                        if len(s) <= 3:
                            continue
                        # 包含中文或字母
                        has_cjk = re.search(r'[\u4e00-\u9fff]', s)
                        has_alpha = any(c.isalpha() for c in s)
                        if has_cjk or has_alpha:
                            lines.append(s)
                    text = "\n".join(lines[:500])
                except Exception:
                    pass
            if not text.strip():
                text = f"[DWG: {os.path.basename(file_path)} — 无文字内容]"
        except Exception as e:
            text = f"[DWG: {os.path.basename(file_path)} — 解析失败: {str(e)[:80]}]"

    elif ext in (".stl", ".obj", ".step", ".stp", ".iges", ".igs", ".gltf", ".glb"):
        try:
            import subprocess
            result = subprocess.run(
                ["strings", file_path],
                capture_output=True, text=True, timeout=30
            )
            lines = []
            for l in result.stdout.split("\n"):
                s = l.strip()
                if len(s) <= 3:
                    continue
                has_cjk = re.search(r'[\u4e00-\u9fff]', s)
                has_alpha = any(c.isalpha() for c in s)
                if has_cjk or has_alpha:
                    lines.append(s)
            # 加上文件元信息
            file_size = os.path.getsize(file_path)
            size_mb = file_size / (1024 * 1024)
            header = f"[3D: {os.path.basename(file_path)}] [{ext.upper()} | {size_mb:.1f} MB]\n"
            text = header + "\n".join(lines[:500])
            if len(text) == len(header):
                text = header + "[此三维文件无可提取的文字信息]"
        except Exception as e:
            text = f"[3D: {os.path.basename(file_path)} — 解析失败: {str(e)[:80]}]",
    return text


# ============ 方案1: PDF 流式分页提取 ============

def _extract_pdf_streaming(file_path: str, pages_per_batch: int = 20) -> list[str]:
    """
    流式分页提取 PDF 文本，每次处理 pages_per_batch 页。
    [来源] poppler-utils pdftotext 支持 -f/-l 参数指定页码范围
    """
    import subprocess
    # 先用 pdfinfo 获取总页数
    result = subprocess.run(
        ["pdfinfo", file_path],
        capture_output=True, text=True, timeout=15
    )
    total_pages = 0
    for line in result.stdout.split("\n"):
        if line.strip().startswith("Pages:"):
            total_pages = int(line.split(":")[1].strip())
            break

    if total_pages == 0:
        # fallback：整体提取
        r = subprocess.run(["pdftotext", "-layout", "-q", file_path, "-"],
                           capture_output=True, text=True, timeout=60)
        return [r.stdout.strip()]

    batches = []
    for start in range(1, total_pages + 1, pages_per_batch):
        end = min(start + pages_per_batch - 1, total_pages)
        try:
            r = subprocess.run(
                ["pdftotext", "-f", str(start), "-l", str(end), "-layout", "-q", file_path, "-"],
                capture_output=True, text=True, timeout=30
            )
            text = r.stdout.strip()
            if text:
                batches.append(text)
        except subprocess.TimeoutExpired:
            print(f"[WARN] PDF 页 {start}-{end} 提取超时，跳过")
            continue
    return batches


# ============ 方案2: 异步任务队列 ============

# 任务状态存储（内存字典，重启丢失未完成任务）
_task_store: dict[str, dict] = {}
_task_lock = threading.Lock()


def _process_file_async(file_path: str, file_name: str, file_hash: str, ext: str, task_id: str):
    """后台处理：提取 → 分段 → 向量化 → 入库"""
    _upload_semaphore.acquire()
    try:
        _update_task(task_id, "extracting", 5, "正在提取文本…")

        # 方案1: PDF 大文件流式分页提取
        if ext == ".pdf":
            batches = _extract_pdf_streaming(file_path)
        else:
            raw = _extract_text(file_path, ext)
            batches = [raw]

        with _task_lock:
            _task_store[task_id]["total_batches"] = len(batches)
        _update_task(task_id, "processing", 10, f"提取完成，共 {len(batches)} 批")

        embedder = _get_embedder()
        all_data = []

        for bi, batch_text in enumerate(batches):
            clean = _clean_text(batch_text)
            if not clean.strip():
                continue

            pct = 10 + int(80 * (bi + 1) / len(batches))
            _update_task(task_id, "processing", pct, f"处理第 {bi+1}/{len(batches)} 批…")

            # Parent-Child 分片
            pc_chunks = _generate_parent_child_chunks(clean)
            if not pc_chunks:
                continue

            # 向量化 Child
            zero_vec = [0.0] * EMBEDDING_DIM
            ci = [i for i, c in enumerate(pc_chunks) if c["chunk_type"] == "child"]
            ct = [pc_chunks[j]["text"] for j in ci]
            cv = embedder.encode(ct, show_progress_bar=False).tolist() if ct else []
            vm = {}
            for j, v in zip(ci, cv):
                vm[j] = v

            now = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")
            cat = _classify_text(clean)
            for j, pc in enumerate(pc_chunks):
                safe_text, _ = _audit_text(pc["text"])
                all_data.append({
                    "vector": vm.get(j, zero_vec),
                    "file_name": file_name,
                    "file_hash": file_hash,
                    "text": safe_text,
                    "chunk_type": pc["chunk_type"],
                    "parent_id": pc["parent_id"],
                    "chunk_index": pc["chunk_index"],
                    "total_chunks": pc["total_chunks"],
                    "category": cat,
                    "tags": "[]",
                    "trust": "unverified",
                    "audit_note": "",
                    "created_at": now,
                })

        # 入库
        _update_task(task_id, "saving", 95, f"正在入库 {len(all_data)} 条…")
        _faiss_insert(all_data)

        # 移动到正式目录
        final_dir = Path(UPLOAD_DIR) / datetime.now().strftime("%Y-%m")
        final_dir.mkdir(parents=True, exist_ok=True)
        if os.path.exists(file_path):
            shutil.move(str(file_path), str(final_dir / os.path.basename(file_path)))

        _update_task(task_id, "done", 100, "入库完成")
        with _task_lock:
            _task_store[task_id]["chunks"] = len(all_data)

    except Exception as e:
        _update_task(task_id, "error", 0, str(e)[:200])
    finally:
        _upload_semaphore.release()


def _update_task(task_id: str, status: str, progress: int, message: str):
    with _task_lock:
        if task_id in _task_store:
            _task_store[task_id].update({
                "status": status, "progress": progress, "message": message,
                "updated_at": datetime.now(timezone.utc).isoformat()
            })


# ============ 结构感知分片 (P0) ============

def _detect_structure_type(ext: str, text: str) -> str:
    """检测文档结构类型"""
    ext = ext.lower()
    if ext in (".cfg", ".conf", ".ini"):
        return "config"
    if ext in (".json", ".xml"):
        return "structured"
    if ext in (".xlsx", ".xls"):
        return "spreadsheet"
    if ext in (".csv", ".log"):
        if _is_tabular(text):
            return "tabular"
        return "log"
    return "generic"

def _is_tabular(text: str) -> bool:
    """检测是否为表格型文本（CSV/TSV）"""
    lines = [l for l in text.split("\n") if l.strip()]
    if len(lines) < 2:
        return False
    # 前 5 行是否逗号/制表符分隔且列数一致
    col_counts = []
    for l in lines[:5]:
        commas = l.count(",")
        tabs = l.count("\t")
        if commas > tabs:
            col_counts.append(commas + 1)
        elif tabs > 0:
            col_counts.append(tabs + 1)
    return len(col_counts) >= 2 and len(set(col_counts)) == 1 and col_counts[0] >= 2

def _split_config(text: str, file_name: str = "") -> list[str]:
    """
    网络设备配置结构感知分片。
    识别 Cisco/华为/H3C/Ruijie 配置块边界：
    - interface / vlan / acl / route-map / ip pool / user-group 等
    每个配置块独立成 chunk，前置块名作为元信息。
    """
    # 配置块开始标记（行首或 ! 后）
    block_start = re.compile(
        r'^[!\s]*(interface|vlan\s+\d+|acl\s|access-list\s|'
        r'route-map\s|ip\s+(route|pool|prefix-list|access-list)\s|'
        r'user-group\s|ssid-profile\s|aaa-profile\s|'
        r'ap-group\s|ap-system-profile\s|wlan\s|'
        r'line\s|router\s|controller\s|class-map\s|policy-map\s|'
        r'dhcp\s+(server|pool|enable)|'
        r'nat\s|firewall\s|security-policy\s|'
        r'web-ac\s|wlan-ac\s)',
        re.IGNORECASE
    )

    lines = text.split("\n")
    chunks = []
    current = []
    current_header = ""

    for line in lines:
        stripped = line.strip()
        # 跳过纯注释和空行（但保留在块内作为上下文）
        if not stripped and not current:
            continue

        m = block_start.match(stripped)
        if m:
            # 遇到新块，保存当前块
            if current:
                block_text = "\n".join(current).strip()
                if block_text and len(block_text) >= 20:
                    # 如果块过长，在内部二次切分
                    if len(block_text) > CHUNK_SIZE:
                        chunks.extend(_split_long_block(block_text))
                    else:
                        chunks.append(block_text)
            current = [line]
            current_header = m.group(0).strip()
        else:
            if current or stripped:
                current.append(line)

    # 最后一个块
    if current:
        block_text = "\n".join(current).strip()
        if block_text and len(block_text) >= 20:
            if len(block_text) > CHUNK_SIZE:
                chunks.extend(_split_long_block(block_text))
            else:
                chunks.append(block_text)

    return chunks if chunks else [text]

def _split_long_block(text: str) -> list[str]:
    """对超长配置块按行二次切分"""
    lines = text.split("\n")
    chunks = []
    current = ""
    for line in lines:
        if len(current) + len(line) + 1 > CHUNK_SIZE and current:
            chunks.append(current.strip())
            current = line
        else:
            current += "\n" + line if current else line
    if current.strip():
        chunks.append(current.strip())
    return chunks

def _split_spreadsheet(text: str, file_name: str = "") -> list[str]:
    """
    Excel/CSV 表格感知分片。
    保留表头作为每个 chunk 的 schema 上下文，
    数据行按 CHUNK_SIZE 分组。
    """
    lines = [l for l in text.split("\n") if l.strip()]
    if not lines:
        return [text]

    chunks = []
    current_sheet = ""
    header = ""
    data_rows = []

    for line in lines:
        if line.startswith("=== ") and line.endswith(" ==="):
            # 保存上一个 sheet
            if header and data_rows:
                chunks.extend(_pack_table_chunks(current_sheet, header, data_rows))
            current_sheet = line
            header = ""
            data_rows = []
        elif not header:
            header = line  # 第一行是表头
        else:
            data_rows.append(line)

    # 最后一个 sheet
    if header and data_rows:
        chunks.extend(_pack_table_chunks(current_sheet, header, data_rows))

    return chunks if chunks else _segment_text(text)

def _pack_table_chunks(sheet_name: str, header: str, rows: list[str]) -> list[str]:
    """将表格行打包成 chunk，每组带表头上下文"""
    chunks = []
    prefix = f"{sheet_name}\n{header}\n" if sheet_name else f"{header}\n"
    current = prefix

    for row in rows:
        if len(current) + len(row) + 1 > CHUNK_SIZE:
            chunks.append(current.strip())
            current = prefix + row
        else:
            current += "\n" + row

    if current.strip() and current != prefix.strip():
        chunks.append(current.strip())

    return chunks

def _split_structured(text: str, file_name: str = "") -> list[str]:
    """
    JSON/XML 结构感知分片。
    JSON：按顶层 key 或数组元素切分，保留路径作为前缀。
    """
    import json as _json

    try:
        data = _json.loads(text)
    except Exception:
        return _segment_text(text)

    chunks = []

    def _flatten(obj, path=""):
        if isinstance(obj, dict):
            # 将每个顶层 key 作为独立 chunk（带路径前缀）
            for key, val in obj.items():
                key_path = f"{path}.{key}" if path else key
                if isinstance(val, (dict, list)):
                    serialized = _json.dumps(val, ensure_ascii=False, indent=2)
                    if len(serialized) > CHUNK_SIZE:
                        _flatten(val, key_path)
                    else:
                        chunks.append(f"[{key_path}]\n{serialized}")
                else:
                    chunks.append(f"[{key_path}]\n{val}")
        elif isinstance(obj, list):
            # 数组元素按 CHUNK_SIZE 分组
            group = []
            group_len = 0
            prefix = f"[{path}]" if path else ""
            for item in obj:
                item_str = _json.dumps(item, ensure_ascii=False, indent=2) if isinstance(item, (dict, list)) else str(item)
                if group_len + len(item_str) > CHUNK_SIZE and group:
                    chunks.append(f"{prefix}\n" + "\n".join(group))
                    group = [item_str]
                    group_len = len(item_str)
                else:
                    group.append(item_str)
                    group_len += len(item_str)
            if group:
                chunks.append(f"{prefix}\n" + "\n".join(group))

    _flatten(data)
    return chunks if chunks else _segment_text(text)

def _clean_text(raw: str) -> str:
    """清洗文本：去 HTML、去乱码、去多余空行、统一空白"""
    # 去 HTML 标签
    cleaned = re.sub(r"<[^>]+>", " ", raw)
    # 去 URL
    cleaned = re.sub(r"https?://\S+", " ", cleaned)
    # 去 PDF 提取产生的乱码字符（常见 OCR/提取噪声）
    cleaned = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]', ' ', cleaned)
    cleaned = re.sub(r'(?:[\u0400-\u04FF]|[\u0500-\u052F])+', ' ', cleaned)  # 西里尔字母混入
    # 去多余的空白行
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    # 压缩连续空格
    cleaned = re.sub(r"[ \t]+", " ", cleaned)
    # 去掉首尾空白
    cleaned = cleaned.strip()
    return cleaned


def _audit_text(text: str) -> tuple[str, list[str]]:
    """内容审核：检测敏感信息，脱敏处理"""
    flags = []
    for pattern in SENSITIVE_PATTERNS:
        if pattern.search(text):
            flags.append(f"检测到敏感信息匹配: {pattern.pattern[:40]}...")
    # 脱敏
    cleaned = SENSITIVE_PATTERNS[0].sub("[已脱敏]", text)
    for p in SENSITIVE_PATTERNS[1:]:
        cleaned = p.sub("[已脱敏]", cleaned)
    return cleaned, flags


def _segment_text(clean_text: str, file_name: str = "", ext: str = "") -> list[str]:
    """
    统一分片入口。
    根据文档类型选择最佳分片策略，然后对结果施加真正的滑动窗口重叠。
    """
    struct_type = _detect_structure_type(ext, clean_text)

    # 阶段 1：结构感知分片
    if struct_type == "config":
        chunks = _split_config(clean_text, file_name)
    elif struct_type == "spreadsheet" or struct_type == "tabular":
        chunks = _split_spreadsheet(clean_text, file_name)
    elif struct_type == "structured":
        chunks = _split_structured(clean_text, file_name)
    else:
        chunks = _split_generic(clean_text)

    # 阶段 2：滑动窗口重叠
    if CHUNK_OVERLAP > 0 and len(chunks) > 1:
        chunks = _apply_overlap(chunks)

    return chunks


def _split_generic(clean_text: str) -> list[str]:
    """通用文本分片：按段落 + 句子边界"""
    paragraphs = [p.strip() for p in clean_text.split("\n") if p.strip()]
    chunks = []
    current = ""

    for para in paragraphs:
        stripped = para.strip()
        if not stripped:
            continue
        # 跳过纯数字/符号行
        if re.match(r"^[\d\s\W]+$", stripped) and not stripped[0].isalpha() and stripped[0] not in ('!', '#', ';', '/'):
            continue
        if len(stripped) < 10 and not any(c.isalpha() for c in stripped):
            continue

        if len(current) + len(stripped) > CHUNK_SIZE:
            if current:
                chunks.append(current.strip())
            current = stripped
        else:
            current += "\n" + stripped if current else stripped

    if current.strip():
        chunks.append(current.strip())

    # 强制拆分超长块
    final_chunks = []
    for chunk in chunks:
        if len(chunk) <= CHUNK_SIZE:
            final_chunks.append(chunk)
        else:
            sentences = re.split(r"(?<=[。！？.!?\n])", chunk)
            sub = ""
            for sent in sentences:
                if len(sub) + len(sent) > CHUNK_SIZE and sub:
                    final_chunks.append(sub.strip())
                    sub = sent
                else:
                    sub += sent
            if sub.strip():
                final_chunks.append(sub.strip())

    return final_chunks


def _apply_overlap(chunks: list[str]) -> list[str]:
    """
    真正的滑动窗口重叠。
    每个 chunk 末尾的 CHUNK_OVERLAP 字符复制到下一个 chunk 的开头，
    确保跨 chunk 边界的关键信息不丢失。
    """
    if not chunks:
        return chunks
    overlapped = [chunks[0]]
    for i in range(1, len(chunks)):
        prev = chunks[i - 1]
        curr = chunks[i]
        # 取上一块的尾部作为当前块的前缀
        tail_len = min(CHUNK_OVERLAP, len(prev))
        overlap_text = prev[-tail_len:]
        # 只在当前块不以该文本开头时才添加
        if not curr.startswith(overlap_text):
            curr = overlap_text + "\n" + curr
        overlapped.append(curr)
    return overlapped


# ============ B1+B2: 查询预处理 ============

# 中文停用词
_STOP_WORDS = set("的 了 在 是 我 有 和 就 不 人 都 一 一个 上 也 很 到 说 要 去 你 会 着 没有 看 好 自己 这 他 她 它 们 那 些 什么 怎么 哪 吗 吧 呢 啊 但 并且 虽然 因为 所以 如果 可以 需要 可能 或者 被 从 把 向 对 与 及 等 之 为 其 将 已 与 该 这个 那个 如何 怎样 为什么 请问 麻烦 帮 我 查 一下 看看 找".split())

# 同义词/术语映射（网络/IT 领域）
_SYNONYMS = {
    "trunk": ["trunk", "干道", "汇聚", "port link-type trunk"],
    "access": ["access", "接入", "port link-type access"],
    "vlan": ["vlan", "虚拟局域网"],
    "ip": ["ip", "地址", "网段", "子网"],
    "交换机": ["交换机", "switch", "lsw"],
    "路由器": ["路由器", "router", "ar"],
    "网线": ["网线", "双绞线", "utp", "rj45", "布线"],
    "防火墙": ["防火墙", "firewall", "acl", "安全策略"],
    "配置": ["配置", "config", "cfg", "设置"],
    "接口": ["接口", "interface", "端口", "port"],
    "拓扑": ["拓扑", "topology", "架构"],
    "dhcp": ["dhcp", "动态主机配置"],
    "dns": ["dns", "域名", "解析"],
    "nat": ["nat", "地址转换"],
    "vpn": ["vpn", "虚拟专用网"],
    "nps": ["nps", "802.1x", "准入", "认证"],
    "wifi": ["wifi", "无线", "wlan", "ssid"],
    "ap": ["ap", "接入点", "无线ap"],
    "ac": ["ac", "无线控制器"],
    "带宽": ["带宽", "速率", "bandwidth"],
    "采购": ["采购", "采购合同", "报价", "供货"],
    "授权": ["授权", "license", "许可"],
    "oee": ["oee", "设备综合效率", "设备效率"],
}

def _preprocess_query(raw: str) -> str:
    """
    B1: 查询预处理——去停用词 + 中文分词 + 关键词提取
    B2: 同义词扩展

    "帮我查一下 LSW1 的 VLAN 101 trunk 配置"
    → "LSW1 VLAN 101 trunk 配置 干道 汇聚 port link-type trunk"
    """
    # 1. 去除口语化前缀
    cleaned = re.sub(r'^(帮我|麻烦|请问|请|我想|我要|查一下|查查|看看|找找|搜索|检索)\s*', '', raw, flags=re.IGNORECASE)
    cleaned = re.sub(r'(一下|吗|呢|吧|啊|呀)\s*$', '', cleaned)

    # 2. jieba 分词（中文）
    if jieba:
        words = list(jieba.cut(cleaned))
    else:
        words = cleaned.split()

    # 3. 过滤：去停用词、短词、纯数字（保留数字+VLAN IP 形式的组合）
    keywords = []
    for w in words:
        w = w.strip().lower()
        if not w or len(w) < 1:
            continue
        if w in _STOP_WORDS:
            continue
        # 保留英文/数字/中文混合词、纯英文、有意义的中文
        if re.search(r'[a-zA-Z\d]', w) or len(w) >= 2:
            keywords.append(w)

    # 4. 保留原始格式的英文词（如 GE0/0/1 不被分词拆散）
    raw_tokens = re.findall(r'[a-zA-Z]+[\d/][\w/.-]*|\d+\.\d+\.\d+\.\d+', cleaned, re.IGNORECASE)
    for rt in raw_tokens:
        if rt.lower() not in [k.lower() for k in keywords]:
            keywords.append(rt.lower())

    # 5. 同义词扩展
    expanded = list(keywords)
    for kw in keywords:
        for syn_key, syn_list in _SYNONYMS.items():
            if kw in syn_list:
                for s in syn_list:
                    if s.lower() not in [e.lower() for e in expanded]:
                        expanded.append(s.lower())
                break

    result = ' '.join(expanded)
    if not result:
        result = raw.strip()
    return result


# ============ P2: BM25 关键词检索器 ============

class _BM25Retriever:
    """简易 BM25 关键词检索，用于与向量检索混合。"""
    def __init__(self):
        self.k1 = 1.5
        self.b = 0.75
        self.documents: list[str] = []
        self.doc_ids: list[int] = []
        self.avgdl = 0.0
        self.idf_cache: dict[str, float] = {}

    def index(self, docs: list[str], ids: list[int]):
        """构建 BM25 索引（内存）"""
        self.documents = docs
        self.doc_ids = ids
        self.avgdl = sum(len(d.split()) for d in docs) / max(len(docs), 1)
        # 预计算 IDF
        N = len(docs)
        df = {}
        for doc in docs:
            for term in set(doc.lower().split()):
                df[term] = df.get(term, 0) + 1
        self.idf_cache = {t: max(0, (N - c + 0.5) / (c + 0.5)) for t, c in df.items()}

    def search(self, query: str, top_k: int = 20) -> list[tuple[int, float]]:
        """返回 [(doc_id, score), ...]"""
        if not self.documents:
            return []
        query_terms = query.lower().split()
        scores = []
        for i, doc in enumerate(self.documents):
            doc_terms = doc.lower().split()
            dl = len(doc_terms)
            score = 0.0
            for term in query_terms:
                if term not in self.idf_cache:
                    continue
                tf = doc_terms.count(term)
                idf = self.idf_cache[term]
                score += idf * (tf * (self.k1 + 1)) / (tf + self.k1 * (1 - self.b + self.b * dl / self.avgdl))
            if score > 0:
                scores.append((self.doc_ids[i], score))
        scores.sort(key=lambda x: x[1], reverse=True)
        return scores[:top_k]


# 全局 BM25 索引（懒更新）
_bm25: Optional[_BM25Retriever] = None
_bm25_version = 0  # 版本号，用于检测是否需要重建


def _get_bm25() -> Optional[_BM25Retriever]:
    """P2: 获取/重建 BM25 索引（基于 FAISS 元数据）"""
    global _bm25, _bm25_version
    if not BM25_ENABLED:
        return None

    if _index is None or _index.ntotal == 0:
        return None

    current_count = _index.ntotal
    if _bm25 is not None and _bm25_version == current_count:
        return _bm25

    try:
        with _meta_lock:
            docs = [m.get("text", "") for m in _metadata if m.get("text")]
            ids = list(range(len(docs)))
        _bm25 = _BM25Retriever()
        _bm25.index(docs, ids)
        _bm25_version = current_count
        return _bm25
    except Exception:
        return None


# ============ P1: Parent-Child 索引 ============

def _generate_parent_child_chunks(full_text: str) -> list[dict]:
    """
    生成 Parent-Child 双层分片。
    返回 list[dict]，每个 dict 包含 chunk 文本和元信息。

    结构：
    - Parent：粗粒度（PARENT_CHUNK_SIZE，默认 1500 字符），保留文档完整段落
    - Child：细粒度（CHILD_CHUNK_SIZE，默认 400 字符），从 Parent 中切出，用于精准匹配

    检索时用 Child 匹配，返回 Parent 完整上下文。
    """
    # 阶段 1：生成 Parent chunks（自然段落边界）
    parents = _split_to_parents(full_text)

    # 阶段 2：每个 Parent 生成 Child chunks
    all_chunks = []
    for parent_idx, parent_text in enumerate(parents):
        parent_id = f"p{parent_idx}"
        # Parent chunk（不含向量，向量化在上层完成）
        all_chunks.append({
            "text": parent_text,
            "chunk_type": "parent",
            "parent_id": parent_id,
            "chunk_index": parent_idx,
            "total_chunks": len(parents),
        })
        # Child chunks
        children = _split_to_children(parent_text)
        for child_idx, child_text in enumerate(children):
            all_chunks.append({
                "text": child_text,
                "chunk_type": "child",
                "parent_id": parent_id,
                "chunk_index": parent_idx,
                "total_chunks": len(parents),
            })

    return all_chunks


def _split_to_parents(text: str) -> list[str]:
    """
    按自然段落边界切分 Parent chunks。
    优先在双换行符（段落边界）处切分，
    尽量保持 PARENT_CHUNK_SIZE 附近的大小。
    """
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    if not paragraphs:
        paragraphs = [p.strip() for p in text.split("\n") if p.strip()]
    if not paragraphs:
        return [text]

    parents = []
    current = ""
    for para in paragraphs:
        if len(current) + len(para) > PARENT_CHUNK_SIZE and current:
            parents.append(current.strip())
            current = para
        else:
            current += "\n\n" + para if current else para
    if current.strip():
        parents.append(current.strip())
    return parents


def _split_to_children(parent_text: str) -> list[str]:
    """从 Parent 中切出 Children，优先在句边界处切分"""
    if len(parent_text) <= CHILD_CHUNK_SIZE:
        return [parent_text]

    # 先尝试按双换行
    parts = [p.strip() for p in parent_text.split("\n\n") if p.strip()]
    if len(parts) == 1:
        parts = [p.strip() for p in parent_text.split("\n") if p.strip()]
    if len(parts) == 1:
        # 按句子切
        parts = [s.strip() for s in re.split(r"(?<=[。！？.!?])", parent_text) if s.strip()]

    children = []
    current = ""
    for part in parts:
        if len(current) + len(part) > CHILD_CHUNK_SIZE and current:
            children.append(current.strip())
            current = part
        else:
            current += "\n" + part if current else part
    if current.strip():
        children.append(current.strip())

    return children if children else [parent_text]


def _classify_text(text: str) -> str:
    """简单分类：基于关键词"""
    text_lower = text.lower()
    categories = {
        "网络建设": ["vlan", "交换机", "路由器", "子网", "网关", "路由", "端口", "防火墙", "拓扑"],
        "IT资产": ["软件", "授权", "账户", "设备", "IP", "电脑", "用户名", "命名", "密码", "license"],
        "采购合同": ["合同", "报价", "采购", "付款", "交货", "供货", "金额", "发票", "供应商"],
        "自动化": ["CT", "OEE", "产线", "自动化", "CMK", "节拍", "传感器", "PLC"],
        "标准件": ["型号", "规格", "米思米", "线轨", "轴承", "气缸", "标准件"],
        "公司架构": ["部门", "架构", "组织", "职责", "制度", "流程", "人事"],
    }
    scores = {}
    for cat, keywords in categories.items():
        scores[cat] = sum(1 for kw in keywords if kw in text_lower)
    if scores:
        best = max(scores, key=scores.get)
        if scores[best] > 0:
            return best
    return "未分类"


# ============ LIFECYCLE ============
@app.on_event("startup")
async def startup():
    """预热模型和数据库"""
    try:
        global _ready
        _ready = False
        _get_embedder()
        dim = _get_embedder().get_embedding_dimension()
        import sys
        this_module = sys.modules[__name__]
        setattr(this_module, 'EMBEDDING_DIM', dim)
        _load_index()
        if not _check_dim_compat():
            print(f"[WARN] 索引维度不匹配 (当前 {dim})，重建…")
            _reset_index()
        _ready = True
        print(f"[OK] 嵌入模型 {EMBEDDING_MODEL} 已加载 (维度={dim})")
        print(f"[OK] FAISS 索引已就绪 @ {DB_PATH}")
        _start_backup_scheduler()
        print(f"[OK] 每日备份调度器已启动（凌晨2点）")
    except Exception as e:
        print(f"[WARN] 启动延迟初始化失败: {e}")


# ============ MIDDLEWARE ============
@app.middleware("http")
async def add_request_id(request: Request, call_next):
    """为每个请求添加唯一 ID，写入 response header"""
    req_id = request.headers.get("X-Request-ID", uuid.uuid4().hex[:12])
    response = await call_next(request)
    response.headers["X-Request-ID"] = req_id
    return response

@app.middleware("http")
async def log_requests(request: Request, call_next):
    """记录所有请求"""
    req_id = request.headers.get("X-Request-ID", "-")
    start = time.time()
    response = await call_next(request)
    elapsed = time.time() - start
    print(f"[{datetime.now(timezone.utc).strftime('%H:%M:%S')}] "
          f"[{req_id}] {request.method} {request.url.path} → {response.status_code} ({elapsed:.3f}s)")
    return response


# ============ ERROR HANDLERS ============
@app.exception_handler(Exception)
async def global_exception_handler(request: Request, exc: Exception):
    return JSONResponse(
        status_code=500,
        content={"error": "Internal server error", "detail": str(exc)[:200]},
    )


# ============ RATE LIMITER ============
_RATE_LIMIT_WINDOW = 60  # 秒
_RATE_LIMIT_MAX = 120    # 每窗口最大请求数
_rate_limit_store: dict = {}  # IP -> (window_start, count)
_rate_lock = threading.Lock()

@app.middleware("http")
async def rate_limiter(request: Request, call_next):
    """简单基于 IP 的速率限制，只对 API 和搜索路径生效"""
    path = request.url.path
    if not path.startswith("/api/"):
        return await call_next(request)
    ip = request.client.host if request.client else "unknown"
    now = time.time()
    with _rate_lock:
        entry = _rate_limit_store.get(ip)
        if entry is None or now - entry[0] > _RATE_LIMIT_WINDOW:
            _rate_limit_store[ip] = (now, 1)
        else:
            count = entry[1] + 1
            _rate_limit_store[ip] = (entry[0], count)
            if count > _RATE_LIMIT_MAX:
                return JSONResponse(status_code=429, content={"error": "请求太频繁，请稍后重试", "retry_after": _RATE_LIMIT_WINDOW})
        # 定期清理过期条目（放在锁外执行减少锁持有时间）
        if len(_rate_limit_store) > 10000:
            stale = [k for k, v in _rate_limit_store.items() if now - v[0] > _RATE_LIMIT_WINDOW * 2]
            for k in stale:
                try:
                    del _rate_limit_store[k]
                except KeyError:
                    pass
        return await call_next(request)


# ============ ROUTES ============
@app.get("/api/health")
async def health():
    """健康检查 — 返回服务状态和基本信息"""
    total_chunks = _index.ntotal if _index else 0
    with _meta_lock:
        total_files = len(set(m.get("file_hash", "") for m in _metadata))
    return {
        "status": "healthy",
        "ready": _ready,
        "total_files": total_files,
        "total_chunks": total_chunks,
        "uptime_seconds": round(time.time() - START_TIME, 1),
    }


def _backup_db():
    """每日备份 FAISS 索引"""
    try:
        _save_index()  # 先落盘
        date_str = datetime.now().strftime("%Y%m%d")
        for fname in ["index.faiss", "metadata.json"]:
            src = os.path.join(DB_PATH, fname)
            if os.path.exists(src):
                shutil.copy2(src, os.path.join(BACKUP_DIR, f"{fname}.{date_str}"))
        # 保留最近 7 天备份
        all_backups = sorted(os.listdir(BACKUP_DIR), reverse=True)
        kept_count = {}  # base_name → 已保留数量
        for old in all_backups:
            # 保留每个文件的最新 7 份
            base = re.sub(r'\.\d{8}$', '', old)
            cnt = kept_count.get(base, 0)
            if cnt < 7:
                kept_count[base] = cnt + 1
            else:
                try:
                    os.remove(os.path.join(BACKUP_DIR, old))
                except OSError:
                    pass
    except Exception:
        pass


def _log_search(query: str, user_ip: str = "", result_count: int = 0, elapsed_ms: float = 0):
    """记录搜索日志（含结果数和耗时）"""
    try:
        today = datetime.now().strftime("%Y%m%d")
        log_path = os.path.join(SEARCH_LOG_DIR, f"search_{today}.jsonl")
        entry = json.dumps({
            "time": datetime.now().isoformat(),
            "query": query[:200],
            "ip": user_ip,
            "results": result_count,
            "elapsed_ms": round(elapsed_ms, 1),
        }, ensure_ascii=False)
        with open(log_path, "a", encoding="utf-8") as f:
            f.write(entry + "\n")
    except Exception:
        pass


def _get_search_history(days: int = 7) -> list[dict]:
    """获取最近 N 天搜索历史，去重取最新 30 条"""
    queries = []
    try:
        now = datetime.now()
        from datetime import timedelta
        for i in range(days):
            day = (now - timedelta(days=i)).strftime("%Y%m%d")
            log_path = os.path.join(SEARCH_LOG_DIR, f"search_{day}.jsonl")
            if os.path.exists(log_path):
                with open(log_path, "r", encoding="utf-8") as f:
                    for line in f:
                        line = line.strip()
                        if line:
                            try:
                                queries.append(json.loads(line))
                            except Exception:
                                pass
    except Exception:
        pass
    # 去重（按 query），取最新 30 条
    seen = set()
    unique = []
    for q in reversed(queries):
        kw = q.get("query", "").strip()
        if kw and kw not in seen:
            seen.add(kw)
            unique.append(kw)
            if len(unique) >= 30:
                break
    return list(reversed(unique))


def _get_hot_queries(days: int = 7, top_n: int = 20) -> list[dict]:
    """获取热门搜索词统计"""
    from collections import Counter
    counter = Counter()
    try:
        now = datetime.now()
        from datetime import timedelta
        for i in range(days):
            day = (now - timedelta(days=i)).strftime("%Y%m%d")
            log_path = os.path.join(SEARCH_LOG_DIR, f"search_{day}.jsonl")
            if os.path.exists(log_path):
                with open(log_path, "r", encoding="utf-8") as f:
                    for line in f:
                        line = line.strip()
                        if line:
                            try:
                                entry = json.loads(line)
                                q = entry.get("query", "").strip()
                                if q:
                                    counter[q] += 1
                            except Exception:
                                pass
    except Exception:
        pass
    return [{"query": q, "count": c} for q, c in counter.most_common(top_n)]


def _start_backup_scheduler():
    """每日凌晨2点备份"""
    def loop():
        while True:
            now = datetime.now()
            next_run = now.replace(hour=2, minute=0, second=0, microsecond=0)
            if now >= next_run:
                next_run = next_run.replace(day=now.day + 1)
            wait = (next_run - now).total_seconds()
            time.sleep(wait)
            _backup_db()
    threading.Thread(target=loop, daemon=True).start()


@app.get("/api/stats", response_model=StatsResponse)
async def stats():
    """获取服务统计"""
    total_chunks = _index.ntotal if _index else 0
    # 快照元数据减少锁时间
    with _meta_lock:
        meta_snapshot = list(_metadata)
    total_files = len(set(m.get("file_hash", "") for m in meta_snapshot))

    # 统计今日反馈数
    fb_count = 0
    try:
        fb_file = os.path.join(SEARCH_LOG_DIR, f"feedback_{datetime.now().strftime('%Y%m%d')}.jsonl")
        if os.path.exists(fb_file):
            with open(fb_file, "r", encoding="utf-8") as f:
                fb_count = sum(1 for line in f if line.strip())
    except Exception:
        pass

    return StatsResponse(
        total_files=total_files,
        total_chunks=total_chunks,
        uptime_seconds=round(time.time() - START_TIME, 1),
        total_tools=len(TOOLS_DATA),
        total_faq=len(FAQ_DATA),
        feedback_count_today=fb_count,
    )


@app.get("/api/documents/{file_hash}")
async def get_document_detail(file_hash: str):
    """获取文档详情（所有 chunk 的完整文本）"""
    with _meta_lock:
        chunks = [m for m in _metadata if m.get("file_hash") == file_hash]
    if not chunks:
        raise HTTPException(status_code=404, detail="文档不存在")
    # 按 parent → child 排序合并
    parents = [c for c in chunks if c.get("chunk_type") == "parent"]
    children = [c for c in chunks if c.get("chunk_type") != "parent"]
    return {
        "file_name": chunks[0].get("file_name", ""),
        "file_hash": file_hash,
        "category": chunks[0].get("category", "未分类"),
        "created_at": chunks[0].get("created_at", ""),
        "parent_count": len(parents),
        "child_count": len(children),
        "total_chunks": len(chunks),
        "preview": [{"text": p.get("text", "")[:800], "chunk_type": p.get("chunk_type")} for p in parents[:5]],
    }


@app.get("/api/documents")
async def list_documents():
    """列出所有已入库文档"""
    with _meta_lock:
        seen = {}
        for m in _metadata:
            fh = m.get("file_hash", "")
            if fh not in seen:
                seen[fh] = {
                    "file_name": m.get("file_name", ""),
                    "file_hash": fh,
                    "category": m.get("category", "未分类"),
                    "created_at": m.get("created_at", ""),
                }
        return {"documents": list(seen.values())}


@app.get("/api/search-history")
async def search_history():
    return {"history": _get_search_history()}


@app.get("/api/search")
async def search(
    q: str = Query(..., min_length=1, max_length=500),
    top_k: int = Query(default=10, ge=1, le=50),
    page: int = Query(default=1, ge=1),
    page_size: int = Query(default=10, ge=1, le=50),
    lang: str = Query(default="auto", regex="^(auto|zh|en)$"),
    request: Request = None,
):
    """
    混合检索：向量语义 + BM25 关键词 → Reranker 精排。
    lang: auto=自动检测, zh=中文优先, en=英文优先
    """
    if not q.strip():
        raise HTTPException(status_code=400, detail="查询内容不能为空")

    # === B1+B2: 查询预处理 ===
    original_q = q.strip()
    processed_q = _preprocess_query(original_q)

    ip = request.client.host if request else ""
    t0 = time.time()

    # 语言检测
    detected_lang = "en" if all(ord(c) < 128 for c in original_q if c.isalpha()) else "zh"

    try:
        if _index is None or _index.ntotal == 0:
            _log_search(original_q, ip, 0, (time.time()-t0)*1000)
            return {"results": [], "query": original_q, "page": 1, "page_size": page_size, "total": 0, "total_pages": 0, "has_more": False, "lang": detected_lang}

        embedder = _get_embedder()
        query_vec = embedder.encode(processed_q).astype(np.float32).reshape(1, -1)
        faiss.normalize_L2(query_vec)  # IP → COSINE

        # === 阶段 1: FAISS 向量检索 (Child chunks only) ===
        # 优化：先快照元数据，减少锁持有时间
        with _meta_lock:
            meta_snapshot = list(_metadata)

        child_indices = [i for i, m in enumerate(meta_snapshot) if m.get("chunk_type") == "child"]

        if not child_indices or _index is None or _index.ntotal == 0:
            _log_search(original_q, ip, 0, (time.time()-t0)*1000)
            return {"results": [], "query": original_q, "page": 1, "page_size": page_size, "total": 0, "total_pages": 0, "has_more": False, "lang": detected_lang}

        # FAISS 全量搜索，后过滤 child（添加超时保护，防止大索引搜索卡死）
        k = min(top_k * 10, _index.ntotal)
        try:
            distances, indices = await asyncio.wait_for(
                asyncio.to_thread(_index.search, query_vec, k),
                timeout=10.0
            )
        except asyncio.TimeoutError:
            _log_search(original_q, ip, 0, (time.time()-t0)*1000)
            raise HTTPException(status_code=504, detail="搜索超时，请尝试缩小查询范围")

        # 过滤 child + 去重
        candidates = {}
        for dist, idx in zip(distances[0], indices[0]):
            if idx < 0 or idx >= len(meta_snapshot):
                continue
            meta = meta_snapshot[idx]
            if meta.get("chunk_type") != "child":
                continue
            parent_id = meta.get("parent_id", "")
            if parent_id not in candidates or dist > candidates[parent_id]["score"]:
                candidates[parent_id] = {"meta": meta, "score": float(dist)}

        # === 阶段 2: BM25 关键词检索 ===
        bm = _get_bm25()
        if bm:
            bm_hits = bm.search(processed_q, top_k=RERANK_TOP_K)
            for meta_idx, bm_score in bm_hits:
                with _meta_lock:
                    if meta_idx >= len(_metadata):
                        continue
                    meta = _metadata[meta_idx]
                if meta.get("chunk_type") != "child":
                    continue
                pid = meta.get("parent_id", "")
                norm_bm = min(bm_score / 100, 0.5)
                if pid in candidates:
                    candidates[pid]["score"] = candidates[pid]["score"] * 0.6 + norm_bm * 0.4
                else:
                    candidates[pid] = {"meta": meta, "score": norm_bm}

        if not candidates:
            _log_search(original_q, ip, 0, (time.time()-t0)*1000)
            return {"results": [], "query": original_q, "page": 1, "page_size": page_size, "total": 0, "total_pages": 0, "has_more": False, "lang": detected_lang}

        # === 阶段 3: Reranker 精排（如有） ===
        reranker = _get_reranker()
        if reranker:
            cand_items = list(candidates.items())
            pairs = [(processed_q, c["meta"].get("text", "")[:512]) for _, c in cand_items]
            try:
                rerank_scores = reranker.predict(pairs, show_progress_bar=False)
                if hasattr(rerank_scores, 'tolist'):
                    rerank_scores = rerank_scores.tolist()
                for i, s in enumerate(rerank_scores):
                    candidates[cand_items[i][0]]["score"] = float(s)
            except Exception as e:
                print(f"[WARN] Reranker 失败: {e}")

        # === 阶段 4: 分数归一化 → 排序 → 回取 Parent ===
        # FAISS IP 距离可能为负，归一化到 0-1 区间（越高越好）
        if candidates:
            scores_list = [c["score"] for c in candidates.values()]
            min_s, max_s = min(scores_list), max(scores_list)
            if max_s > min_s:
                for c in candidates.values():
                    c["score"] = (c["score"] - min_s) / (max_s - min_s)
            else:
                for c in candidates.values():
                    c["score"] = 0.5

        sorted_cands = sorted(candidates.items(), key=lambda x: x[1]["score"], reverse=True)

        # 收集 parent_id，回取 parent 文本
        parent_ids = [pid for pid, _ in sorted_cands[:top_k * 5]]
        parent_cache = {}
        # 在已快照的元数据中查找 parent
        for m in meta_snapshot:
            if m.get("chunk_type") == "parent" and m.get("parent_id") in parent_ids:
                pid = m["parent_id"]
                if pid not in parent_cache:
                    parent_cache[pid] = m

        # 组装输出
        output = []
        seen_files = set()
        for pid, cdata in sorted_cands:
            meta = cdata["meta"]
            file_key = meta.get("file_hash", "") + meta.get("file_name", "")
            if file_key in seen_files:
                continue
            seen_files.add(file_key)

            parent = parent_cache.get(pid)
            display_text = parent.get("text", "") if parent else meta.get("text", "")
            display_category = parent.get("category", meta.get("category", "未分类")) if parent else meta.get("category", "未分类")

            output.append({
                "file_name": meta.get("file_name", ""),
                "text_preview": (display_text or "")[:500],
                "text": display_text or "",
                "category": display_category,
                "tags": json.loads(parent.get("tags", "[]") or "[]") if parent else json.loads(meta.get("tags", "[]") or "[]"),
                "trust": parent.get("trust", "unverified") if parent else meta.get("trust", "unverified"),
                "audit_note": parent.get("audit_note", "") if parent else meta.get("audit_note", ""),
                "score": round(cdata["score"], 4),
            })

            if len(output) >= top_k:
                break

        # 分页：先取全部 top_k 结果，再按 page/page_size 切片
        total_count = len(output)
        start = (page - 1) * page_size
        paged = output[start:start + page_size]
        total_pages = max(1, (total_count + page_size - 1) // page_size) if total_count > 0 else 0

        elapsed = (time.time() - t0) * 1000
        _log_search(original_q, ip, len(paged), elapsed)
        return {
            "results": paged,
            "query": original_q,
            "elapsed_ms": round(elapsed, 1),
            "page": page,
            "page_size": page_size,
            "total": total_count,
            "total_pages": total_pages,
            "has_more": page < total_pages,
            "lang": detected_lang,
        }
    except Exception as e:
        elapsed = (time.time() - t0) * 1000
        _log_search(original_q, ip, 0, elapsed)
        raise HTTPException(status_code=500, detail=str(e)[:200])


@app.post("/api/upload/preview", response_model=PreviewResponse)
async def upload_preview(file: UploadFile = File(...)):
    """
    上传文件并预览分段结果（不立即入库）。
    前端展示分段后，用户确认后才真正入库。
    """
    # 验证扩展名
    ext = Path(file.filename or "unknown").suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(status_code=400, detail=f"不支持的文件格式: {ext}，支持: {', '.join(sorted(ALLOWED_EXTENSIONS))}")

    safe_name = _sanitize_filename(file.filename or "untitled")
    content = await file.read()

    # 大小检查
    if len(content) > MAX_FILE_MB * 1024 * 1024:
        raise HTTPException(status_code=400, detail=f"文件过大，最大 {MAX_FILE_MB}MB")
    if ext == ".pdf" and len(content) > MAX_PDF_MB * 1024 * 1024:
        raise HTTPException(status_code=400, detail=f"PDF 最大 {MAX_PDF_MB}MB，当前 {len(content)/1024/1024:.1f}MB")

    # 计算哈希
    file_hash = hashlib.sha256(content).hexdigest()

    # 保存到临时目录
    temp_dir = Path(UPLOAD_DIR) / "temp"
    temp_dir.mkdir(parents=True, exist_ok=True)
    temp_path = temp_dir / f"{file_hash}_{safe_name}"
    temp_path.write_bytes(content)

    # 提取文本
    try:
        raw_text = _extract_text(str(temp_path), ext)
    except Exception as e:
        temp_path.unlink(missing_ok=True)
        raise HTTPException(status_code=500, detail=f"文件解析失败: {e}")

    # 清洗
    clean_text = _clean_text(raw_text)

    # 分段
    chunks = _segment_text(clean_text, safe_name, ext)

    if not chunks:
        temp_path.unlink(missing_ok=True)
        raise HTTPException(status_code=400, detail=f"文档内容为空（提取到 {len(clean_text)} 字符但未能分段）")

    # 审核每个分段
    preview_chunks = []
    for i, chunk in enumerate(chunks):
        safe_chunk, flags = _audit_text(chunk)
        preview_chunks.append(ChunkPreview(
            index=i,
            text=safe_chunk,
            length=len(safe_chunk),
            flagged=bool(flags),
            flag_reason="; ".join(flags) if flags else "",
        ))

    return PreviewResponse(
        file_name=safe_name,
        file_hash=file_hash,
        total_chunks=len(preview_chunks),
        chunks=preview_chunks,
    )


@app.post("/api/upload/confirm", response_model=UploadResponse)
async def upload_confirm(body: ConfirmRequest):
    """
    确认入库：用户编辑分段后，向量化并存入 Milvus。
    """
    if not body.chunks:
        raise HTTPException(status_code=400, detail="分段列表不能为空")

    embedder = _get_embedder()

    # 检查是否已入库（FAISS 元数据查重）
    with _meta_lock:
        existing = any(m.get("file_hash") == body.file_hash for m in _metadata)
    if existing:
        # 删除旧数据
        with _meta_lock:
            _metadata[:] = [m for m in _metadata if m.get("file_hash") != body.file_hash]
        _rebuild_faiss_from_meta()

    # 分类
    full_text = "\n".join(c["text"] for c in body.chunks)
    category = _classify_text(full_text)

    # 审核
    all_tags = set()
    for chunk in body.chunks:
        safe_text, _ = _audit_text(chunk["text"])
        chunk["text"] = safe_text
        all_tags.update(chunk.get("tags", []))

    # 批量向量化 & 插入（Parent-Child 双层，仅向量化 Child）
    pc_chunks = _generate_parent_child_chunks(full_text)
    # 分离 parent 和 child，只向量化 child（parent 存零向量占位）
    zero_vec = [0.0] * EMBEDDING_DIM
    child_indices = [i for i, c in enumerate(pc_chunks) if c["chunk_type"] == "child"]
    child_texts = [pc_chunks[i]["text"] for i in child_indices]
    child_vectors = embedder.encode(child_texts, show_progress_bar=False).tolist() if child_texts else []
    vec_map = {}
    for cidx, vec in zip(child_indices, child_vectors):
        vec_map[cidx] = vec

    now_iso = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")
    data = []
    for i, pc in enumerate(pc_chunks):
        v = vec_map.get(i, zero_vec)
        data.append({
            "vector": v,
            "file_name": body.file_name,
            "file_hash": body.file_hash,
            "text": pc["text"],
            "chunk_type": pc["chunk_type"],
            "parent_id": pc["parent_id"],
            "chunk_index": pc["chunk_index"],
            "total_chunks": pc["total_chunks"],
            "category": category,
            "tags": json.dumps(sorted(all_tags), ensure_ascii=False),
            "trust": "unverified",
            "audit_note": "",
            "created_at": now_iso,
        })

    # FAISS 插入
    _faiss_insert(data)

    # 移动文件到最终目录
    temp_path = Path(UPLOAD_DIR) / "temp" / f"{body.file_hash}_{body.file_name}"
    final_dir = Path(UPLOAD_DIR) / datetime.now().strftime("%Y-%m")
    final_dir.mkdir(parents=True, exist_ok=True)
    if temp_path.exists():
        shutil.move(str(temp_path), str(final_dir / body.file_name))

    return UploadResponse(
        file=body.file_name,
        chunks=len(body.chunks),
        category=category,
    )


@app.post("/api/upload", response_model=UploadResponse)
async def upload_simple(file: UploadFile = File(...)):
    """
    上传文件入库（方案2: 异步处理）。
    文件 ≤ 5MB 同步处理，> 5MB 后台异步。
    """
    ext = Path(file.filename or "unknown").suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(status_code=400, detail=f"不支持的文件格式: {ext}")

    safe_name = _sanitize_filename(file.filename or "untitled")
    content = await file.read()
    file_size_mb = len(content) / (1024 * 1024)

    if len(content) > MAX_FILE_MB * 1024 * 1024:
        raise HTTPException(status_code=400, detail=f"文件过大，最大 {MAX_FILE_MB}MB")

    file_hash = hashlib.sha256(content).hexdigest()

    # 查重
    with _meta_lock:
        if any(m.get("file_hash") == file_hash for m in _metadata):
            return UploadResponse(file=safe_name, chunks=0, category="已存在", status="duplicate")

    # 保存文件
    temp_dir = Path(UPLOAD_DIR) / "temp"
    temp_dir.mkdir(parents=True, exist_ok=True)
    temp_path = temp_dir / f"{file_hash}_{safe_name}"
    temp_path.write_bytes(content)

    # 小文件同步处理，大文件异步
    if file_size_mb <= 5:
        return _upload_sync(temp_path, safe_name, file_hash, ext)
    else:
        return _upload_async(temp_path, safe_name, file_hash, ext)


def _upload_sync(temp_path, safe_name, file_hash, ext):
    """同步处理：小文件直接入库"""
    raw_text = _extract_text(str(temp_path), ext)
    clean_text = _clean_text(raw_text)

    if not clean_text.strip():
        temp_path.unlink(missing_ok=True)
        raise HTTPException(status_code=400, detail=f"文档内容为空")

    embedder = _get_embedder()
    pc_chunks = _generate_parent_child_chunks(clean_text)

    if not pc_chunks:
        temp_path.unlink(missing_ok=True)
        raise HTTPException(status_code=400, detail=f"文档内容为空（提取到 {len(clean_text)} 字符但未能分段）")

    zero_vec = [0.0] * EMBEDDING_DIM
    ci = [i for i, c in enumerate(pc_chunks) if c["chunk_type"] == "child"]
    ct = [pc_chunks[j]["text"] for j in ci]
    cv = embedder.encode(ct, show_progress_bar=False).tolist() if ct else []
    vm = {}
    for j, v in zip(ci, cv):
        vm[j] = v

    now_iso = datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ")
    cat = _classify_text(clean_text)
    data = []
    for i, pc in enumerate(pc_chunks):
        safe_text, _ = _audit_text(pc["text"])
        data.append({
            "vector": vm.get(i, zero_vec),
            "file_name": safe_name,
            "file_hash": file_hash,
            "text": safe_text,
            "chunk_type": pc["chunk_type"],
            "parent_id": pc["parent_id"],
            "chunk_index": pc["chunk_index"],
            "total_chunks": pc["total_chunks"],
            "category": cat,
            "tags": "[]",
            "trust": "unverified",
            "audit_note": "",
            "created_at": now_iso,
        })

    _faiss_insert(data)

    final_dir = Path(UPLOAD_DIR) / datetime.now().strftime("%Y-%m")
    final_dir.mkdir(parents=True, exist_ok=True)
    if temp_path.exists():
        shutil.move(str(temp_path), str(final_dir / safe_name))

    return UploadResponse(file=safe_name, chunks=len(data), category=cat)


def _upload_async(temp_path, safe_name, file_hash, ext):
    """异步处理：大文件后台处理，立即返回任务 ID"""
    task_id = uuid.uuid4().hex[:12]
    with _task_lock:
        _task_store[task_id] = {
            "task_id": task_id,
            "file_name": safe_name,
            "file_hash": file_hash,
            "status": "queued",
            "progress": 0,
            "message": "已加入处理队列",
            "chunks": 0,
            "total_batches": 0,
            "created_at": datetime.now(timezone.utc).isoformat(),
        }
    t = threading.Thread(target=_process_file_async,
                         args=(str(temp_path), safe_name, file_hash, ext, task_id),
                         daemon=True)
    t.start()
    return UploadResponse(file=safe_name, chunks=0, category=f"处理中 ({task_id})", status="processing")


@app.get("/api/task/{task_id}")
async def get_task_status(task_id: str):
    """查询异步任务状态"""
    try:
        with _task_lock:
            task = _task_store.get(task_id)
        if not task:
            raise HTTPException(status_code=404, detail="任务不存在")
        return task
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"查询任务状态失败: {str(e)[:200]}")


@app.delete("/api/documents/{file_hash}")
async def delete_document(file_hash: str):
    """删除已入库的文档"""
    global _index, _metadata
    with _meta_lock:
        before = len(_metadata)
        _metadata = [m for m in _metadata if m.get("file_hash") != file_hash]
        after = len(_metadata)
        deleted = before - after
    if deleted == 0:
        raise HTTPException(status_code=404, detail="文档不存在")
    # 重建 FAISS 索引
    _rebuild_faiss_from_meta()
    return {"deleted": deleted, "file_hash": file_hash}


# ============ ADMIN API ============
# 管理面板专用 API

@app.get("/api/admin/stats")
async def admin_stats():
    """
    获取综合统计数据
    返回：文件数、片段数、按格式分类、数据库大小
    """
    total_chunks = _index.ntotal if _index else 0
    format_dist = {}
    total_files = 0
    with _meta_lock:
        meta_snapshot = list(_metadata)
    seen = set()
    for m in meta_snapshot:
        fh = m.get("file_hash", "")
        if fh not in seen:
            seen.add(fh)
            total_files += 1
            ext = Path(m.get("file_name", "")).suffix.lower().lstrip('.') or "unknown"
            format_dist[ext] = format_dist.get(ext, 0) + 1

    # 计算数据库大小
    db_size = 0
    if os.path.exists(DB_PATH):
        for root, dirs, files in os.walk(DB_PATH):
            db_size += sum(os.path.getsize(os.path.join(root, f)) for f in files)

    return {
        "totalFiles": total_files,
        "totalChunks": total_chunks,
        "dbSize": db_size,
        "formatDistribution": format_dist,
    }


@app.get("/api/admin/upload-trend")
async def admin_upload_trend():
    """
    获取每日上传趋势
    返回：最近 7 天每天的文件上传数量
    """
    try:
        from datetime import timedelta
        trend = []
        now = datetime.now()
        with _meta_lock:
            for i in range(6, -1, -1):
                day = now - timedelta(days=i)
                day_str = day.strftime("%Y-%m-%d")
                count = sum(1 for m in _metadata if m.get("created_at", "").startswith(day_str))
                trend.append({
                    "date": day.strftime("%m-%d"),
                    "count": count,
                })
        
        return trend
    except Exception as e:
        print(f"获取上传趋势失败：{e}")
        return []


@app.get("/api/admin/error-logs")
async def admin_error_logs():
    """
    获取最近的错误日志
    读取：~/kb-server/logs/error.log
    """
    try:
        log_path = os.path.expanduser("~/kb-server/logs/error.log")
        logs = []
        
        if os.path.exists(log_path):
            with open(log_path, "r", encoding="utf-8", errors="replace") as f:
                lines = f.readlines()[-100]  # 最近 100 行
                for line in reversed(lines):
                    line = line.strip()
                    if not line:
                        continue
                    # 尝试解析日志格式
                    try:
                        # 假设 JSON 格式
                        entry = json.loads(line)
                        logs.append(entry)
                    except:
                        # 非 JSON 格式，尝试解析文本日志
                        if "ERROR" in line or "error" in line or "Exception" in line:
                            logs.append({
                                "level": "ERROR",
                                "message": line[:500],
                                "timestamp": datetime.now().isoformat(),
                            })
                    if len(logs) >= 50:  # 最多返回 50 条
                        break
        
        return logs
    except Exception as e:
        print(f"读取错误日志失败：{e}")
        return []


@app.get("/api/admin/ai-search-logs")
async def admin_ai_search_logs():
    """
    获取 AI 搜索记录
    读取：~/kb-server/logs/search_*.jsonl
    """
    try:
        from datetime import timedelta
        logs = []
        now = datetime.now()
        
        # 读取最近 7 天的搜索日志
        for i in range(7):
            day = now - timedelta(days=i)
            day_str = day.strftime("%Y%m%d")
            log_path = os.path.expanduser(f"~/kb-server/logs/search_{day_str}.jsonl")
            
            if os.path.exists(log_path):
                with open(log_path, "r", encoding="utf-8") as f:
                    for line in f:
                        line = line.strip()
                        if line:
                            try:
                                entry = json.loads(line)
                                # 转换为前端格式
                                logs.append({
                                    "timestamp": entry.get("time", entry.get("timestamp")),
                                    "query": entry.get("query", ""),
                                    "userId": entry.get("ip", "anonymous"),
                                    "status": "success",
                                    "responseTime": 0,
                                })
                            except:
                                pass
        
        # 按时间倒序，最多返回 200 条
        logs.sort(key=lambda x: x.get("timestamp", ""), reverse=True)
        return logs[:200]
    except Exception as e:
        print(f"读取搜索日志失败：{e}")
        return []


@app.get("/api/admin/server-status")
async def admin_server_status():
    """
    获取服务器状态
    返回：CPU、内存、磁盘使用情况
    注意：CPU 采样通过 run_in_executor 异步执行，不阻塞事件循环。
    """
    try:
        import psutil
        import concurrent.futures

        loop = asyncio.get_event_loop()

        # CPU 采样在独立线程中执行（interval=1 会阻塞 1 秒）
        cpu_percent = await loop.run_in_executor(
            concurrent.futures.ThreadPoolExecutor(max_workers=1),
            lambda: psutil.cpu_percent(interval=1)
        )

        # 内存使用
        memory = psutil.virtual_memory()
        memory_percent = memory.percent
        memory_total = memory.total

        # 磁盘使用（根分区）
        disk = psutil.disk_usage("/")
        disk_percent = disk.percent
        disk_total = disk.total

        # 运行时间
        uptime_seconds = time.time() - START_TIME
        days = int(uptime_seconds // 86400)
        hours = int((uptime_seconds % 86400) // 3600)
        uptime_str = f"{days}天 {hours}小时"

        return {
            "cpu": cpu_percent,
            "memory": memory_percent,
            "memoryTotal": memory_total,
            "disk": disk_percent,
            "diskTotal": disk_total,
            "uptime": uptime_str,
        }
    except ImportError:
        # psutil 未安装，返回模拟数据
        print("警告：psutil 未安装，无法获取服务器状态")
        return {
            "cpu": 0,
            "memory": 0,
            "memoryTotal": 0,
            "disk": 0,
            "diskTotal": 0,
            "uptime": "未知",
        }
    except Exception as e:
        print(f"获取服务器状态失败：{e}")
        return {
            "cpu": 0,
            "memory": 0,
            "memoryTotal": 0,
            "disk": 0,
            "diskTotal": 0,
            "uptime": "错误",
        }


@app.get("/api/admin/recent-activities")
async def admin_recent_activities():
    """
    获取最近操作记录
    返回：最近的文件上传、删除等操作
    """
    try:
        with _meta_lock:
            seen = {}
            for m in sorted(_metadata, key=lambda x: x.get("created_at", ""), reverse=True):
                fh = m.get("file_hash", "")
                if fh not in seen:
                    seen[fh] = m
            activities = []
            for r in list(seen.values())[:50]:
                activities.append({
                    "type": "upload",
                    "action": "upload",
                    "filename": r.get("file_name", "未知文件"),
                    "fileType": Path(r.get("file_name", "")).suffix.lstrip('.') or "unknown",
                    "status": "completed",
                    "timestamp": r.get("created_at"),
                    "progress": 100,
                })
        return activities
    except Exception as e:
        print(f"获取活动记录失败：{e}")
        return []


# ============ 前端新 API（阶段六：后端代码补充） ============

@app.get("/api/tools")
async def get_tools():
    """获取工具列表 — 替代前端硬编码 TOOL_URLS，优先读动态配置"""
    cfg = _load_config()
    tools = cfg.get("tools", TOOLS_DATA)
    return {"tools": tools, "total": len(tools)}


@app.get("/api/tools/check")
async def check_tools_availability():
    """检测所有工具 URL 的可达性（异步并发）"""
    import asyncio, concurrent.futures
    cfg = _load_config()
    tools = cfg.get("tools", TOOLS_DATA)

    async def check_one(tool):
        url = tool.get("url", "")
        if not url or url.startswith("#"):
            return {"id": tool.get("id"), "name": tool.get("name"), "available": tool.get("available", True), "reachable": None, "latency_ms": None, "error": "未配置 URL"}
        if url.startswith("\\\\"):
            return {"id": tool.get("id"), "name": tool.get("name"), "available": tool.get("available", True), "reachable": None, "latency_ms": None, "error": "SMB 路径，跳过检测"}
        try:
            import urllib.request
            t0 = time.time()
            req = urllib.request.Request(url, method="HEAD")
            loop = asyncio.get_event_loop()
            resp = await loop.run_in_executor(
                concurrent.futures.ThreadPoolExecutor(max_workers=1),
                lambda: urllib.request.urlopen(req, timeout=5)
            )
            elapsed = (time.time() - t0) * 1000
            return {"id": tool.get("id"), "name": tool.get("name"), "reachable": True, "latency_ms": round(elapsed, 1), "status_code": resp.status}
        except Exception as e:
            return {"id": tool.get("id"), "name": tool.get("name"), "reachable": False, "latency_ms": None, "error": str(e)[:100]}

    tasks = [check_one(t) for t in tools]
    results = await asyncio.gather(*tasks)
    online = sum(1 for r in results if r.get("reachable") is True)
    return {"results": results, "online": online, "offline": len(results) - online, "total": len(results)}


@app.get("/api/faq")
async def get_faq():
    """获取常见问题列表 — 替代前端硬编码 QUESTIONS，优先读动态配置"""
    cfg = _load_config()
    faq = cfg.get("faq", FAQ_DATA)
    return {"faq": faq, "total": len(faq)}


@app.post("/api/feedback")
async def submit_feedback(req: FeedbackRequest):
    """收集 AI 回答反馈 — 替代前端 alert()"""
    try:
        today = datetime.now().strftime("%Y%m%d")
        fb_path = os.path.join(SEARCH_LOG_DIR, f"feedback_{today}.jsonl")
        entry = json.dumps({
            "time": req.timestamp or datetime.now().isoformat(),
            "query": req.query[:500],
            "answer_preview": req.answer_preview[:500],
            "useful": req.useful,
        }, ensure_ascii=False)
        with open(fb_path, "a", encoding="utf-8") as f:
            f.write(entry + "\n")
        return {"status": "ok"}
    except Exception as e:
        print(f"记录反馈失败：{e}")
        return {"status": "error", "message": str(e)[:200]}


# ============ 数据持久化辅助函数 ============
_CONFIG_PATH = os.path.join(DB_PATH, "config.json")
_CONFIG_HISTORY_DIR = os.path.join(DB_PATH, "config_history")

def _load_config() -> dict:
    """加载可编辑配置（工具+FAQ）"""
    try:
        if os.path.exists(_CONFIG_PATH):
            with open(_CONFIG_PATH, "r", encoding="utf-8") as f:
                return json.load(f)
    except Exception:
        pass
    return {"tools": TOOLS_DATA, "faq": FAQ_DATA}

def _save_config(config: dict):
    """保存配置到磁盘，同时保留历史版本（v3.5: 增加磁盘满/权限异常保护）"""
    try:
        os.makedirs(DB_PATH, exist_ok=True)
        # 如果旧配置存在且不同，先归档
        if os.path.exists(_CONFIG_PATH):
            try:
                with open(_CONFIG_PATH, "r", encoding="utf-8") as f:
                    old = f.read()
                new = json.dumps(config, ensure_ascii=False, indent=2)
                if old != new:
                    os.makedirs(_CONFIG_HISTORY_DIR, exist_ok=True)
                    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
                    hist_path = os.path.join(_CONFIG_HISTORY_DIR, f"config_{ts}.json")
                    with open(hist_path, "w", encoding="utf-8") as f:
                        f.write(old)
                    # 保留最近 50 个历史版本
                    try:
                        hfiles = sorted(os.listdir(_CONFIG_HISTORY_DIR), reverse=True)
                        for fname in hfiles[50:]:
                            os.remove(os.path.join(_CONFIG_HISTORY_DIR, fname))
                    except Exception:
                        pass  # 清理失败不影响主流程
            except Exception:
                pass
        # 原子写入：先写临时文件再 rename
        tmp_path = _CONFIG_PATH + ".tmp"
        with open(tmp_path, "w", encoding="utf-8") as f:
            json.dump(config, f, ensure_ascii=False, indent=2)
        os.replace(tmp_path, _CONFIG_PATH)  # 原子替换
    except Exception as e:
        print(f"[ERROR] 保存配置失败: {e}")
        # 清理临时文件
        try:
            tmp = _CONFIG_PATH + ".tmp"
            if os.path.exists(tmp):
                os.remove(tmp)
        except Exception:
            pass


# ============ 管理面板扩展 API ============

@app.get("/api/admin/tools")
async def admin_get_tools():
    """管理面板：获取工具列表（含可编辑版本）"""
    cfg = _load_config()
    return {"tools": cfg.get("tools", TOOLS_DATA), "total": len(cfg.get("tools", TOOLS_DATA))}


@app.post("/api/admin/tools")
async def admin_save_tools(request: Request):
    """管理面板：保存工具列表"""
    try:
        body = await request.json()
        cfg = _load_config()
        cfg["tools"] = body.get("tools", TOOLS_DATA)
        _save_config(cfg)
        return {"status": "ok", "total": len(cfg["tools"])}
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e)[:200])


@app.get("/api/admin/faq")
async def admin_get_faq():
    """管理面板：获取 FAQ 列表"""
    cfg = _load_config()
    return {"faq": cfg.get("faq", FAQ_DATA), "total": len(cfg.get("faq", FAQ_DATA))}


@app.post("/api/admin/faq")
async def admin_save_faq(request: Request):
    """管理面板：保存 FAQ 列表"""
    try:
        body = await request.json()
        cfg = _load_config()
        cfg["faq"] = body.get("faq", FAQ_DATA)
        _save_config(cfg)
        return {"status": "ok", "total": len(cfg["faq"])}
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e)[:200])


@app.get("/api/admin/feedbacks")
async def admin_get_feedbacks(days: int = Query(7, ge=1, le=90)):
    """管理面板：获取最近 N 天的 AI 反馈数据"""
    feedbacks = []
    try:
        now = datetime.now()
        from datetime import timedelta
        for i in range(days):
            day = (now - timedelta(days=i)).strftime("%Y%m%d")
            fb_path = os.path.join(SEARCH_LOG_DIR, f"feedback_{day}.jsonl")
            if os.path.exists(fb_path):
                with open(fb_path, "r", encoding="utf-8") as f:
                    for line in f:
                        line = line.strip()
                        if line:
                            try:
                                feedbacks.append(json.loads(line))
                            except Exception:
                                pass
    except Exception:
        pass
    # 按时间倒序
    feedbacks.sort(key=lambda x: x.get("time", ""), reverse=True)
    total = len(feedbacks)
    useful = sum(1 for fb in feedbacks if fb.get("useful"))
    return {
        "feedbacks": feedbacks[:200],
        "total": total,
        "useful": useful,
        "useless": total - useful,
        "useful_rate": round(useful / total * 100, 1) if total > 0 else 0,
    }


@app.get("/api/admin/hot-queries")
async def admin_hot_queries(days: int = Query(7, ge=1, le=90), top_n: int = Query(20, ge=5, le=100)):
    """管理面板：热门搜索词统计"""
    hot = _get_hot_queries(days, top_n)
    return {"hot_queries": hot, "total_unique": len(hot), "days": days}


@app.get("/api/admin/export/documents")
async def admin_export_documents():
    """导出文档列表为 CSV"""
    from fastapi.responses import StreamingResponse
    import io
    with _meta_lock:
        seen = {}
        for m in _metadata:
            fh = m.get("file_hash", "")
            if fh not in seen:
                seen[fh] = m
    output = io.StringIO()
    output.write("file_name,category,chunks,created_at\n")
    for m in seen.values():
        name = (m.get("file_name", "") or "").replace('"', '""')
        cat = (m.get("category", "") or "").replace('"', '""')
        chunks = sum(1 for x in _metadata if x.get("file_hash") == m.get("file_hash", ""))
        ts = m.get("created_at", "")
        output.write(f'"{name}","{cat}",{chunks},"{ts}"\n')
    csv = output.getvalue()
    output.close()
    return StreamingResponse(io.BytesIO(csv.encode('utf-8-sig')),
                             media_type="text/csv",
                             headers={"Content-Disposition": f"attachment; filename=documents_{datetime.now().strftime('%Y%m%d')}.csv"})


@app.get("/api/admin/export/search-logs")
async def admin_export_search_logs(days: int = Query(7, ge=1, le=90)):
    """导出搜索日志为 CSV"""
    from fastapi.responses import StreamingResponse
    import io
    logs = []
    try:
        now = datetime.now()
        from datetime import timedelta
        for i in range(days):
            day = (now - timedelta(days=i)).strftime("%Y%m%d")
            log_path = os.path.join(SEARCH_LOG_DIR, f"search_{day}.jsonl")
            if os.path.exists(log_path):
                with open(log_path, "r", encoding="utf-8") as f:
                    for line in f:
                        line = line.strip()
                        if line:
                            try:
                                logs.append(json.loads(line))
                            except Exception:
                                pass
    except Exception:
        pass
    logs.sort(key=lambda x: x.get("time", ""), reverse=True)
    output = io.StringIO()
    output.write("time,query,results,elapsed_ms,ip\n")
    for l in logs[:5000]:
        t = (l.get("time", "") or "").replace('"', '""')
        q = (l.get("query", "") or "").replace('"', '""')
        r = l.get("results", "")
        e = l.get("elapsed_ms", "")
        ip = l.get("ip", "")
        output.write(f'"{t}","{q}",{r},{e},"{ip}"\n')
    csv = output.getvalue()
    output.close()
    return StreamingResponse(io.BytesIO(csv.encode('utf-8-sig')),
                             media_type="text/csv",
                             headers={"Content-Disposition": f"attachment; filename=search_logs_{datetime.now().strftime('%Y%m%d')}.csv"})


@app.get("/api/admin/config")
async def admin_get_config():
    """管理面板：获取完整配置（工具+FAQ）"""
    cfg = _load_config()
    return cfg


@app.get("/api/admin/config/history")
async def admin_config_history():
    """管理面板：获取配置变更历史列表"""
    history = []
    try:
        if os.path.exists(_CONFIG_HISTORY_DIR):
            for fname in sorted(os.listdir(_CONFIG_HISTORY_DIR), reverse=True):
                if fname.startswith("config_") and fname.endswith(".json"):
                    ts_str = fname.replace("config_", "").replace(".json", "")
                    fpath = os.path.join(_CONFIG_HISTORY_DIR, fname)
                    try:
                        with open(fpath, "r", encoding="utf-8") as f:
                            data = json.load(f)
                        history.append({
                            "version": ts_str,
                            "tools_count": len(data.get("tools", [])),
                            "faq_count": len(data.get("faq", [])),
                            "file": fname,
                        })
                    except Exception:
                        pass
    except Exception:
        pass
    return {"history": history, "total": len(history)}


@app.post("/api/admin/config/rollback")
async def admin_config_rollback(request: Request):
    """管理面板：回滚配置到指定历史版本"""
    try:
        body = await request.json()
        version_file = body.get("version", "")
        if not version_file or ".." in version_file or "/" in version_file:
            raise HTTPException(status_code=400, detail="无效的版本号")
        hist_path = os.path.join(_CONFIG_HISTORY_DIR, version_file)
        if not os.path.exists(hist_path):
            raise HTTPException(status_code=404, detail="历史版本不存在")
        with open(hist_path, "r", encoding="utf-8") as f:
            rolled = json.load(f)
        _save_config(rolled)
        return {"status": "ok", "version": version_file,
                "tools": len(rolled.get("tools", [])),
                "faq": len(rolled.get("faq", []))}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)[:200])


# ============ MAIN ============
@app.post("/api/reset")
async def reset_collection(request: Request):
    """清空知识库所有数据（需要管理令牌）"""
    token = request.headers.get("x-admin-token", "")
    expected = os.getenv("KB_ADMIN_TOKEN", "polygon-admin-2024")
    if token != expected:
        raise HTTPException(status_code=403, detail="需要管理令牌")
    try:
        count = _index.ntotal if _index else 0
        _reset_index()
        return {"message": f"已清空 {count} 条数据", "deleted": count}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)[:200])


if __name__ == "__main__":
    import uvicorn
    max_body = int(os.getenv("KB_MAX_FILE_MB", "200")) * 1024 * 1024
    uvicorn.run(app, host=HOST, port=PORT, log_level="info", limit_max_requests=10000, timeout_keep_alive=300, limit_concurrency=50)
